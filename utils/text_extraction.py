import io
import fitz  # PyMuPDF
import docx
import pandas as pd

def extract_text_from_file(filename, file_bytes):
    extension = filename.split(".")[-1].lower()
    if extension == "pdf":
        return extract_pdf_text(file_bytes)
    elif extension == "docx":
        return extract_docx_text(file_bytes)
    elif extension in ["xls", "xlsx"]:
        return extract_excel_text(file_bytes)
    else:
        return ""

def extract_pdf_text(file_bytes):
    text = ""
    with fitz.open("pdf", file_bytes) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_docx_text(file_bytes):
    text = ""
    file_stream = io.BytesIO(file_bytes)
    doc = docx.Document(file_stream)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_excel_text(file_bytes):
    file_stream = io.BytesIO(file_bytes)
    df = pd.read_excel(file_stream, engine='openpyxl')
    return df.to_string()




import os
import boto3
import tempfile
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
import pptx
import openpyxl
import textract

AWS S3 Client Setup
s3 = boto3.client(
"s3",
aws_access_key_id=os.getenv("AWS_ACCESS_KEY"),
aws_secret_access_key=os.getenv("AWS_SECRET_KEY"),
region_name="ap-south-1"
)

BUCKET_NAME = "bmc-automation"
FOLDER_NAME = "NAAC_documents"

def download_file_from_s3(object_key: str) -> BytesIO:
"""
Downloads a file from S3 and returns it as a BytesIO object.
"""
response = s3.get_object(Bucket=BUCKET_NAME, Key=object_key)
return BytesIO(response["Body"].read())

def extract_text_from_file(object_key: str) -> str:
file_stream = download_file_from_s3(object_key)
file_ext = object_key.lower().split('.')[-1]

if file_ext == "pdf":
    reader = PdfReader(file_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


elif file_ext == "docx":
    document = Document(file_stream)
    return "\n".join([para.text for para in document.paragraphs])

elif file_ext == "pptx":
    presentation = pptx.Presentation(file_stream)
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

elif file_ext == "xlsx":
    wb = openpyxl.load_workbook(file_stream, data_only=True)
    extracted = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            extracted.append(" ".join([str(cell) if cell else "" for cell in row]))
    return "\n".join(extracted)

elif file_ext == "txt":
    return file_stream.read().decode("utf-8")

else:
    # fallback using textract for uncommon formats
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
        tmp.write(file_stream.read())
        tmp_path = tmp.name
    try:
        return textract.process(tmp_path).decode("utf-8")
    finally:
        os.remove(tmp_path)